from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG = RGBColor(0xFD, 0xF9, 0xF8)
PINK = RGBColor(0xA0, 0x70, 0x6A)
LIGHT_PINK = RGBColor(0xC4, 0x92, 0x8A)
DARK = RGBColor(0x3D, 0x30, 0x32)
GRAY = RGBColor(0x8C, 0x7A, 0x7C)
LIGHT_GRAY = RGBColor(0xB8, 0xA8, 0xA9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MINT = RGBColor(0xB8, 0xD0, 0xC2)
TAG_BG = RGBColor(0xF5, 0xEC, 0xE9)

IMGDIR = "D:/portfolio/images"

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, font_name='Microsoft YaHei', alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_img(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        if width and height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top))

def add_tag_row(slide, left, top, tags):
    x = left
    for t in tags:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(len(t)*0.18+0.3), Inches(0.32))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TAG_BG
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = t
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = PINK
        tf.paragraphs[0].font.name = 'Microsoft YaHei'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        x += len(t)*0.18 + 0.4

# ===== Slide 1: Cover =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Faded lily fills entire slide as background
add_img(slide, f"{IMGDIR}/hero-bg-faded.png", 0, 0, 13.333, 7.5)

add_text_box(slide, 1.5, 1.8, 8, 1.5, '罗楚婷', 60, DARK, True, 'Microsoft YaHei')
add_text_box(slide, 1.5, 3.4, 8, 0.8, '新媒体运营 · 作品集', 28, GRAY, False, 'Microsoft YaHei')

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.3), Inches(1.5), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = PINK
line.line.fill.background()

add_text_box(slide, 1.5, 4.6, 8, 0.5, 'AI-native 运营实践者 · 3000+用户产品创造者 · 14万+曝光', 14, LIGHT_GRAY, False, 'Microsoft YaHei')

# ===== Slide 2: Personal Info =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img(slide, f"{IMGDIR}/avatar.jpg", 0.8, 0.8, 1.5, 1.9)

add_text_box(slide, 2.8, 0.8, 6, 0.7, '个人简介', 36, PINK, True)
add_text_box(slide, 2.8, 1.7, 8, 0.5, '罗楚婷  |  23岁  |  广州  |  2025届  |  (86) 176-7563-7717  |  1822977645@qq.com', 13, GRAY)

# Education
add_text_box(slide, 1, 3.2, 11, 0.5, '教育背景', 18, PINK, True)
add_text_box(slide, 1, 3.8, 11, 0.8, '广东财经大学 · 汉语言文学（商务文秘）· 本科', 16, DARK, True)
add_text_box(slide, 1, 4.3, 11, 0.8, 'GPA 3.84 / 5.0（专业前10%）· 院长奖学金（前1%）· 优秀毕业生 · 英语六级 · 电子商务课程满分', 14, GRAY)

# Skills
add_text_box(slide, 1, 5.2, 11, 0.5, '核心能力', 18, PINK, True)
add_text_box(slide, 1, 5.7, 11, 0.8, '需求调研 · 用户洞察 · PRD撰写 · 功能设计 · 数据复盘 · 用户增长 · 私域运营 · 内容策划\nAI编程(vibe coding) · Procreate · 基础PS · AI生图 · 可灵AI视频生成 · 英语六级', 13, GRAY)

# ===== Slide 3: 目录 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 0.5, 10, 0.7, '目录', 36, PINK, True)

items = [
    ('01', '个人简介', '教育背景 · 核心能力'),
    ('02', 'AI 产品设计与运营', '「菠萝烤盐」考研择校工具  ·  古诗记忆闪卡 PRD'),
    ('03', '内容运营', '小红书个人账号 · 自然流运营与数据复盘'),
    ('04', '活动策划', '暑假班课程策划 · 项目策划与落地'),
    ('05', '视觉设计', 'Procreate 海报设计 · PS · AI 生图'),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.6 + i * 0.9
    add_text_box(slide, 1.5, y, 1, 0.5, num, 28, PINK, True)
    add_text_box(slide, 3, y, 8, 0.45, title, 20, DARK, True)
    add_text_box(slide, 3, y + 0.4, 8, 0.35, desc, 13, GRAY)

# ===== Slide 4: AI产品设计与运营 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 0.5, 10, 0.7, '02  AI 产品设计与运营', 32, PINK, True)

# Left: 菠萝烤盐
add_text_box(slide, 1, 1.3, 5.5, 0.45, '「菠萝烤盐」· 考研择校工具', 20, DARK, True)
add_text_box(slide, 1, 1.7, 5.5, 0.25, '2026.06 – 至今  |  3000+用户 · 14万+曝光', 11, GRAY)
add_text_box(slide, 1, 2.1, 5.8, 0.25, '为什么做', 13, PINK, True)
add_text_box(slide, 1, 2.35, 5.8, 0.65, '市面择校工具只堆信息，用户看完还是不知道"我该报哪个"。考研人真正焦虑的不是"有哪些学校"，是"我能不能考上"。', 11, GRAY)
add_text_box(slide, 1, 3.05, 5.8, 0.25, '做了什么', 13, PINK, True)
add_text_box(slide, 1, 3.3, 5.8, 1.15, '收集134所高校16维数据（含给分松紧、一志愿保护等隐性指标），设计加权评分算法输出冲/稳/保+匹配度评分。vibe coding开发部署，小红书引流+4个粉丝群反馈迭代。', 11, GRAY)
add_text_box(slide, 1, 4.5, 2.8, 0.3, '3000+', 24, PINK, True)
add_text_box(slide, 1, 4.8, 2.8, 0.25, '主动用户', 10, GRAY)
add_text_box(slide, 3.2, 4.5, 2.8, 0.3, '14万+', 24, PINK, True)
add_text_box(slide, 3.2, 4.8, 2.8, 0.25, '笔记曝光 · 双板块第一', 10, GRAY)
add_text_box(slide, 1, 5.2, 5.8, 0.25, '洞察：用户要的不是信息，是判断——"黑犀牛能力"', 11, LIGHT_PINK, True)
add_text_box(slide, 1, 5.5, 5.8, 0.25, 'karka2666.github.io/kaoyan-selecter', 9, GRAY)

# Divider
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.1), Inches(1.3), Inches(0.015), Inches(5.0))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xEC, 0xE0, 0xDE); div.line.fill.background()

# Right: 闪卡
add_text_box(slide, 7.5, 1.3, 5.5, 0.45, '古诗记忆闪卡 · 产品设计', 20, DARK, True)
add_text_box(slide, 7.5, 1.7, 5.5, 0.25, '2026.03 – 2026.05  |  产品设计师', 11, GRAY)
add_text_box(slide, 7.5, 2.1, 5.5, 0.25, '为什么做', 13, PINK, True)
add_text_box(slide, 7.5, 2.35, 5.5, 0.65, '考研背古诗总忘。背单词App有记忆曲线，背诗为什么没有？', 11, GRAY)
add_text_box(slide, 7.5, 3.05, 5.5, 0.25, '做了什么', 13, PINK, True)
add_text_box(slide, 7.5, 3.3, 5.5, 1.15, '访谈考研学生→基于艾宾浩斯曲线设计三种模式（新学/复习/补缺）→设计数据库→页面路由→交互流程→输出完整PRD', 11, GRAY)
add_text_box(slide, 7.5, 4.5, 5.5, 0.25, '产出', 13, PINK, True)
add_text_box(slide, 7.5, 4.75, 5.5, 0.65, '33KB完整PRD（需求→功能→数据→技术方案），产品岗核心作品集', 11, GRAY)
add_text_box(slide, 7.5, 5.5, 5.5, 0.25, 'karka2666.github.io/resume/poetry-flashcard-prd.md', 9, GRAY)

# ===== Slide 5: 内容运营 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 0.5, 10, 0.7, '03  内容运营', 32, PINK, True)
add_text_box(slide, 1, 1.05, 10, 0.35, '小红书个人账号  |  2025.12 – 至今  |  1500+粉丝 · 纯自然流', 13, GRAY)

add_text_box(slide, 1, 1.9, 11, 0.4, '▎为什么做', 16, DARK, True)
add_text_box(slide, 1, 2.3, 11, 0.5, '想验证自己能不能纯靠内容感觉跑自然流量，不做垂直赛道。', 13, GRAY)

add_text_box(slide, 1, 3.0, 11, 0.4, '▎做了什么', 16, DARK, True)
add_text_box(slide, 1, 3.4, 11, 0.8, '选题、文案、封面全自己来，精心挑选生活感和利他性封面建立辨识度。借助AI辅助创作+每次发布复盘数据快速调整方向。非垂直账号稳定踩中爆款，靠的是对平台内容节奏的敏感度和快速试错迭代。', 13, GRAY)

# Show XHS images
for i in range(4):
    path = f"{IMGDIR}/xhs-0{i+1}.jpg"
    if os.path.exists(path):
        add_img(slide, path, 1 + i*3, 4.5, 2.7, 2.2)

add_text_box(slide, 1, 6.9, 11, 0.3, '产出：一套可复用的内容运营方法（选题敏感度→AI辅助创作→数据复盘→快速迭代）', 12, PINK, True)

# ===== Slide 6: 暑假班 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 0.5, 10, 0.7, '04  活动策划', 32, PINK, True)
add_text_box(slide, 1, 1.05, 10, 0.35, '暑假班课程策划  |  2025.06 – 2026.06  |  报名超预期20% · 满意度100%', 13, GRAY)

add_text_box(slide, 1, 1.9, 11, 0.4, '▎为什么做', 16, DARK, True)
add_text_box(slide, 1, 2.3, 11, 0.5, '在教培机构工作时发现暑假班方案空缺，主动承接。', 13, GRAY)

add_text_box(slide, 1, 3.0, 11, 0.4, '▎做了什么', 16, DARK, True)
add_text_box(slide, 1, 3.4, 11, 0.8, '访谈家长提炼"提分效率"和"时间灵活性"两大需求→设计基础班/冲刺班两档差异化产品→分别定价→协调师资教务→按期开课。', 13, GRAY)

add_text_box(slide, 1, 4.4, 11, 0.4, '▎效果与产出', 16, DARK, True)
add_text_box(slide, 1, 4.9, 2, 0.5, '超预期20%', 30, PINK, True)
add_text_box(slide, 1, 5.4, 2, 0.4, '报名人数', 12, GRAY)
add_text_box(slide, 4, 4.9, 2, 0.5, '100%', 30, PINK, True)
add_text_box(slide, 4, 5.4, 2, 0.4, '家长满意度', 12, GRAY)
add_text_box(slide, 1, 6.0, 11, 0.5, '产出：PRD + 运营方案 + PPT + 排版PDF全套方案文档，方案获机构直接采纳执行', 13, GRAY)

# ===== Slide 7: Design Gallery =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 0.5, 10, 0.7, '05  视觉设计', 32, PINK, True)
add_text_box(slide, 1, 1.05, 10, 0.35, 'Procreate 海报设计 · 基础 PS · AI 生图', 13, GRAY)

# Compact 4x3 grid (11 images + 1 empty)
cols, rows = 4, 3
img_w, img_h = 2.7, 2.0
start_x, start_y = 0.8, 1.5
gap_x, gap_y = 0.15, 0.2
for i in range(11):
    path = f"{IMGDIR}/design-{i+1:02d}.jpg"
    col = i % cols
    row = i // cols
    if os.path.exists(path):
        add_img(slide, path, start_x + col*(img_w+gap_x), start_y + row*(img_h+gap_y), img_w, img_h)

# ===== Slide 8: End =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_img(slide, f"{IMGDIR}/hero-bg-faded.png", 0, 0, 13.333, 7.5)

add_text_box(slide, 2, 2, 6, 1.5, '感谢观看', 60, PINK, True)
add_text_box(slide, 2, 3.8, 6, 0.5, '— 罗楚婷 · 2026 —', 20, GRAY)
add_text_box(slide, 2, 5.0, 6, 0.5, '(86) 176-7563-7717  ·  1822977645@qq.com  ·  广州', 13, LIGHT_GRAY)
add_text_box(slide, 2, 5.4, 6, 0.5, 'karka2666.github.io/resume', 12, GRAY)

# Save
out = "D:/罗楚婷-作品集-v6.pptx"
prs.save(out)
print(f"PPT generated: {out}")
print(f"Slides: {len(prs.slides)}")
